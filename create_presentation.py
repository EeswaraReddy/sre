"""
AI, GenAI, and Agentic Systems Presentation Generator
Converts markdown content to PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    TITLE_COLOR = RGBColor(26, 35, 126)  # Deep blue
    ACCENT_COLOR = RGBColor(255, 87, 34)  # Orange
    TEXT_COLOR = RGBColor(33, 33, 33)  # Dark gray
    BG_COLOR = RGBColor(255, 255, 255)  # White
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title = title_frame.paragraphs[0]
    title.text = "AI, Generative AI & Agentic Systems"
    title.font.size = Pt(54)
    title.font.bold = True
    title.font.color.rgb = TITLE_COLOR
    title.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle = subtitle_frame.paragraphs[0]
    subtitle.text = "Framework Comparison & Evaluation Strategies"
    subtitle.font.size = Pt(32)
    subtitle.font.color.rgb = ACCENT_COLOR
    subtitle.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date = date_frame.paragraphs[0]
    date.text = "A Comprehensive Technical Overview - 2026"
    date.font.size = Pt(20)
    date.font.color.rgb = TEXT_COLOR
    date.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Traditional AI vs GenAI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "What is AI and Generative AI?"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Traditional AI (Discriminative AI)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    for text in ["Rule-based expert systems", "Classification and prediction", 
                 "Pattern recognition from existing data", "Examples: Spam filters, fraud detection"]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Generative AI (GenAI)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    for text in ["Creates new content across modalities (text, images, video, audio)",
                 "Foundation models trained on massive datasets",
                 "Transformers and large-scale neural networks",
                 "Examples: GPT-4, Claude, Gemini, DALL-E"]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 3: 2026 Key Trends
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "2026 Key Trends in AI"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    trends = [
        ("✨ Multimodal Integration", "Seamless text, image, audio, video processing"),
        ("🤖 Agentic Capabilities", "Proactive, autonomous task completion"),
        ("🎯 Domain Specialization", "Industry-specific fine-tuned models"),
        ("⚡ Edge AI", "On-device inference for privacy and speed"),
        ("📊 Hyper-personalization", "Real-time adaptive user experiences"),
        ("🔄 Full Process Automation", "End-to-end workflow execution"),
        ("🏢 Enterprise Integration", "Deep embedding in critical business processes")
    ]
    
    for trend_title, trend_desc in trends:
        p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.text = trend_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        p = tf.add_paragraph()
        p.text = trend_desc
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 4: Understanding Agentic AI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Understanding Agentic AI Systems"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "AI Agents are autonomous systems that can:"
    p.font.size = Pt(22)
    p.font.bold = True
    
    capabilities = [
        "🎯 Reason - Analyze situations and make decisions",
        "📋 Plan - Break down complex tasks into steps",
        "🔧 Act - Execute tasks using available tools",
        "🔄 Learn - Improve from feedback and experience",
        "🤝 Collaborate - Work with other agents and humans"
    ]
    
    for cap in capabilities:
        p = tf.add_paragraph()
        p.text = cap
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Key Components"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    components = [
        "Foundation Model - LLM powering reasoning",
        "Tools & APIs - Database, web search, external services",
        "Memory System - Short-term and long-term storage",
        "Environment - Task queue and feedback loops"
    ]
    
    for comp in components:
        p = tf.add_paragraph()
        p.text = comp
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 5: Multi-Agent Systems
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Multi-Agent Systems & Orchestration"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "2026 Trends in Multi-Agent Systems"
    p.font.size = Pt(24)
    p.font.bold = True
    
    mas_points = [
        "Orchestration - Coordinating specialized agents for complex workflows",
        "Collaboration Patterns - Swarm, hierarchical, workflow-based",
        "Role Specialization - Each agent has specific expertise",
        "Prediction: 70% of MAS will feature narrow, focused roles by 2027",
        "Governance-First Design - Permission boundaries, audit logs",
        "Human-in-the-Loop - Approval checkpoints for critical decisions"
    ]
    
    for point in mas_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 6: AI Evaluation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AI Evaluation Frameworks"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "2026 Shift: From 'Can AI do this?' to 'How well, at what cost, and for whom?'"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "Three Evaluation Pillars"
    p.font.size = Pt(24)
    p.font.bold = True
    
    pillars = [
        ("Performance Metrics", ["Accuracy/F1 Score", "BLEU/ROUGE", "Relevance Score"]),
        ("Safety & Ethics", ["Bias Detection", "Toxicity Check", "Fairness Score"]),
        ("Production Quality", ["Hallucination Rate", "Latency/Cost", "Compliance"])
    ]
    
    for pillar_name, metrics in pillars:
        p = tf.add_paragraph()
        p.text = pillar_name
        p.level = 1
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        for metric in metrics:
            p = tf.add_paragraph()
            p.text = metric
            p.level = 2
            p.font.size = Pt(16)
    
    # Slide 7: Evaluation Platforms
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Leading Evaluation Platforms (2026)"
    
    # Add table
    rows, cols = 8, 3
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(4.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(4)
    table.columns[2].width = Inches(3)
    
    # Header row
    headers = ["Platform", "Strengths", "Best For"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    data = [
        ["DeepEval", "Developer-focused, RAG metrics", "Development & Testing"],
        ["Galileo AI", "Hallucination detection", "Production GenAI"],
        ["Arize", "ML observability, drift", "Enterprise Monitoring"],
        ["Patronus AI", "Rubric-based scoring", "Structured Evaluation"],
        ["MLflow", "Experiment tracking", "Custom Workflows"],
        ["RAGAS", "RAG-specific evaluation", "RAG Systems"],
        ["Braintrust", "Dev workflow integration", "End-to-End Platform"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Slide 8: Strands SDK
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Strands SDK Framework"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Open-source, code-first Python framework by AWS for production-ready AI agents"
    p.font.size = Pt(18)
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "Key Features"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    features = [
        "🎯 Simplicity - Minimal boilerplate, opinionated design",
        "🔄 Model Agnostic - Works with Bedrock, OpenAI, Anthropic, Google",
        "🏗️ Production-Ready - Built-in observability, tracing, deployment",
        "🤖 Multi-Agent Support - Swarm, Graph, Workflow patterns",
        "⚡ AWS Integration - Seamless with Bedrock, Lambda, Step Functions"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Use Cases"
    p.font.size = Pt(22)
    p.font.bold = True
    
    for use_case in ["Autonomous incident resolution (SRE)", 
                     "Multi-step workflow automation",
                     "Research and analysis agents"]:
        p = tf.add_paragraph()
        p.text = use_case
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 9: AWS AgentCore
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AWS Bedrock AgentCore Platform"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Agentic platform for building, deploying, and operating AI agents securely at scale"
    p.font.size = Pt(18)
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "Core Services"
    p.font.size = Pt(24)
    p.font.bold = True
    
    services = [
        ("Runtime Service", "Serverless, scalable execution environment"),
        ("Memory Service", "Short-term and long-term context management"),
        ("Gateway Service", "Secure connections to tools and resources"),
        ("Observability", "Production monitoring and quality tracking"),
        ("Policy Engine", "Controls agent-to-tool interactions")
    ]
    
    for service_name, service_desc in services:
        p = tf.add_paragraph()
        p.text = service_name
        p.level = 1
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        p = tf.add_paragraph()
        p.text = service_desc
        p.level = 2
        p.font.size = Pt(14)
    
    # Slide 10: Framework Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Agent Framework Comparison (2026)"
    
    # Add comparison table
    rows, cols = 5, 5
    left = Inches(0.3)
    top = Inches(1.8)
    width = Inches(9.4)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2)
    table.columns[4].width = Inches(1.4)
    
    # Header row
    headers = ["Framework", "Philosophy", "Architecture", "Best For", "Learning Curve"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    data = [
        ["LangChain/\nLangGraph", "Modular, flexible", "Graph-based state machines", "Complex enterprise RAG", "Moderate-High"],
        ["CrewAI", "Role-based teams", "Hierarchical/sequential", "Content creation, SOPs", "Easy"],
        ["AutoGen", "Multi-agent conversation", "Dialogue orchestration", "Code gen, research", "Moderate"],
        ["Strands SDK", "Code-first simplicity", "Lightweight agent loop", "AWS-native production", "Easy-Moderate"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Slide 11: LLM Models Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Latest LLM Models Comparison (2026)"
    
    # Add model comparison table
    rows, cols = 7, 5
    left = Inches(0.2)
    top = Inches(1.8)
    width = Inches(9.6)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    headers = ["Vendor", "Model", "Strengths", "Context Window", "Key Features"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    data = [
        ["OpenAI", "GPT-5.2", "Reasoning, coding", "Extended", "Super-assistant"],
        ["Anthropic", "Claude Opus 4.6", "Agentic coding, safety", "1M tokens", "Computer use"],
        ["Google", "Gemini 3 Pro", "Multimodal, integration", "2M tokens", "Deep Research"],
        ["AWS", "Amazon Nova", "Enterprise, cost-effective", "Variable", "Bedrock-native"],
        ["Meta", "Llama 3+", "Open-source", "Variable", "On-prem deployment"],
        ["Microsoft", "Copilot (GPT-4)", "Productivity", "Extended", "M365 integration"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Slide 12: Cloud Services
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Cloud AI Services Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    providers = [
        ("Amazon Web Services (AWS)", [
            "Amazon Bedrock - Multi-model platform (Claude, Nova, Llama)",
            "SageMaker - Custom ML model training and deployment",
            "Amazon Q - AI assistant for development",
            "AgentCore - Agentic platform for secure agent operations"
        ]),
        ("Microsoft Azure", [
            "Azure OpenAI - Enterprise GPT-4 access",
            "Azure ML - Custom model development",
            "Microsoft Copilot - M365 deep integration",
            "Semantic Kernel - LLM integration framework"
        ]),
        ("Google Cloud Platform (GCP)", [
            "Vertex AI - Unified ML platform",
            "Gemini API - Multimodal AI models",
            "Duet AI - Workspace integration",
            "Agent Builder - Agentic workflow creation"
        ])
    ]
    
    for provider_name, services in providers:
        p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.text = provider_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        for service in services:
            p = tf.add_paragraph()
            p.text = service
            p.level = 1
            p.font.size = Pt(14)
    
    # Slide 13: Use Cases
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Real-World Enterprise Use Cases"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    use_cases = [
        ("Customer Service & Support", [
            "Intelligent ticket routing and resolution",
            "24/7 automated support with escalation"
        ]),
        ("DevOps & SRE", [
            "Incident detection and auto-remediation",
            "Log analysis and root cause identification"
        ]),
        ("Data & Analytics", [
            "Automated data pipeline monitoring",
            "Natural language to SQL query generation"
        ]),
        ("Software Development", [
            "Code generation and review",
            "Bug detection and automated fixing"
        ]),
        ("Finance & Compliance", [
            "Regulatory compliance checking",
            "Fraud detection and risk assessment"
        ])
    ]
    
    for uc_title, uc_details in use_cases:
        p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.text = uc_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        for detail in uc_details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
            p.font.size = Pt(14)
    
    # Slide 14: Best Practices
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Best Practices & Recommendations"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    practices = [
        ("1. Start with Clear Scope", [
            "Define specific tasks and success criteria",
            "Establish human oversight requirements"
        ]),
        ("2. Choose the Right Framework", [
            "AWS-native → Strands SDK + AgentCore",
            "Complex Workflow → LangGraph",
            "Quick Prototype → CrewAI"
        ]),
        ("3. Implement Robust Evaluation", [
            "Pre-production testing with representative data",
            "Continuous monitoring in production",
            "Hallucination and safety checks"
        ]),
        ("4. Governance & Security", [
            "Define permission boundaries",
            "Implement audit logging",
            "Set approval workflows for critical actions"
        ]),
        ("5. Iterate Based on Metrics", [
            "Track task success rates",
            "Monitor tool usage patterns",
            "Continuous prompt optimization"
        ])
    ]
    
    for practice_title, practice_details in practices:
        p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.text = practice_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        for detail in practice_details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
            p.font.size = Pt(12)
    
    # Slide 15: Summary & Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary & Next Steps"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Key Takeaways"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    takeaways = [
        "✅ AI agents are moving from experimental to enterprise-critical",
        "✅ Multi-agent orchestration is the future of complex automation",
        "✅ Evaluation and observability are non-negotiable",
        "✅ Choose frameworks based on use case, not hype",
        "✅ Security and governance must be designed in from day one"
    ]
    
    for takeaway in takeaways:
        p = tf.add_paragraph()
        p.text = takeaway
        p.level = 1
        p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "Recommended Next Steps"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    next_steps = [
        "1. Pilot Project - Start with a well-defined use case",
        "2. Framework Selection - Match framework to requirements",
        "3. Evaluation Setup - Implement monitoring from day one",
        "4. Team Training - Upskill on chosen framework",
        "5. Production Deployment - Iterative rollout"
    ]
    
    for step in next_steps:
        p = tf.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(14)
    
    # Save presentation
    prs.save('e:/Antigravity/sreagent/AI_GenAI_Agents_Presentation.pptx')
    print("Presentation created successfully!")
    return prs

if __name__ == "__main__":
    create_presentation()
