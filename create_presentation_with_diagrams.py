"""
AI, GenAI, and Agentic Systems Presentation Generator with Diagrams
Enhanced version with visual diagrams for each slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

def add_connector(slide, x1, y1, x2, y2, color=RGBColor(100, 100, 100)):
    """Add a connector line between two points"""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

def add_shape_with_text(slide, shape_type, left, top, width, height, text, bg_color, text_color=RGBColor(255, 255, 255), font_size=14):
    """Add a shape with centered text"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color
    
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(font_size)
    text_frame.paragraphs[0].font.color.rgb = text_color
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return shape

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    TITLE_COLOR = RGBColor(26, 35, 126)  # Deep blue
    ACCENT_COLOR = RGBColor(255, 87, 34)  # Orange
    TEXT_COLOR = RGBColor(33, 33, 33)  # Dark gray
    BG_COLOR = RGBColor(255, 255, 255)  # White
    BLUE_1 = RGBColor(33, 150, 243)  # Light blue
    BLUE_2 = RGBColor(3, 169, 244)  # Cyan blue
    GREEN = RGBColor(76, 175, 80)  # Green
    PURPLE = RGBColor(156, 39, 176)  # Purple
    ORANGE = RGBColor(255, 152, 0)  # Orange
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add decorative circles
    for i, (x, y, color) in enumerate([(1, 1, BLUE_1), (8.5, 1, ACCENT_COLOR), (1, 6, GREEN), (8.5, 6, PURPLE)]):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = color
        circle.fill.transparency = 0.3
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title = title_frame.paragraphs[0]
    title.text = "AI, Generative AI & Agentic Systems"
    title.font.size = Pt(54)
    title.font.bold = True
    title.font.color.rgb = TITLE_COLOR
    title.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle = subtitle_frame.paragraphs[0]
    subtitle.text = "Framework Comparison & Evaluation Strategies"
    subtitle.font.size = Pt(32)
    subtitle.font.color.rgb = ACCENT_COLOR
    subtitle.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date = date_frame.paragraphs[0]
    date.text = "A Comprehensive Technical Overview - 2026"
    date.font.size = Pt(20)
    date.font.color.rgb = TEXT_COLOR
    date.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Traditional AI vs GenAI with Evolution Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "What is AI and Generative AI?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Evolution diagram
    box_width = 2
    box_height = 0.8
    y_pos = 1.5
    
    # Traditional AI
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, y_pos, box_width, box_height, 
                       "Traditional AI", BLUE_1, font_size=16)
    
    # Arrow
    add_connector(slide, 0.5 + box_width, y_pos + 0.4, 3.5, y_pos + 0.4, TITLE_COLOR)
    
    # Machine Learning
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.5, y_pos, box_width, box_height, 
                       "Machine\nLearning", BLUE_2, font_size=16)
    
    # Arrow
    add_connector(slide, 3.5 + box_width, y_pos + 0.4, 6.5, y_pos + 0.4, TITLE_COLOR)
    
    # Generative AI
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.5, y_pos, box_width, box_height, 
                       "Generative AI", GREEN, font_size=16)
    
    # Content boxes below
    content_y = 2.8
    
    # Traditional AI details
    trad_box = slide.shapes.add_textbox(Inches(0.3), Inches(content_y), Inches(3), Inches(2.5))
    tf = trad_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Traditional AI"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_1
    
    for text in ["• Rule-based systems", "• Classification & prediction", "• Pattern recognition", "• Spam filters, fraud detection"]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    
    # GenAI details
    gen_box = slide.shapes.add_textbox(Inches(3.5), Inches(content_y), Inches(3), Inches(2.5))
    tf = gen_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Generative AI"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    for text in ["• Creates new content", "• Multimodal (text, image, video)", "• Foundation models", "• GPT-4, Claude, Gemini"]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    
    # 2026 Trends
    trends_box = slide.shapes.add_textbox(Inches(6.8), Inches(content_y), Inches(2.8), Inches(2.5))
    tf = trends_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2026 Trends"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    for text in ["✨ Multimodal", "🤖 Agentic AI", "🎯 Specialized", "⚡ Edge AI"]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(13)
        p.space_after = Pt(4)
    
    # Slide 3: 2026 Key Trends with Icon Grid
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Key Trends in AI"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Trend boxes in grid
    trends = [
        ("✨", "Multimodal Integration", "Text, image, audio, video", BLUE_1),
        ("🤖", "Agentic Capabilities", "Autonomous task completion", GREEN),
        ("🎯", "Domain Specialization", "Industry-specific models", PURPLE),
        ("⚡", "Edge AI", "On-device inference", ORANGE),
        ("📊", "Hyper-personalization", "Real-time adaptation", BLUE_2),
        ("🔄", "Full Automation", "End-to-end processes", RGBColor(233, 30, 99)),
    ]
    
    col = 0
    row = 0
    start_x = 0.5
    start_y = 1.5
    box_width = 3
    box_height = 1.8
    gap = 0.3
    
    for icon, title, desc, color in trends:
        x = start_x + col * (box_width + gap)
        y = start_y + row * (box_height + gap)
        
        # Create rounded rectangle
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(x), Inches(y), Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.fill.transparency = 0.2
        box.line.color.rgb = color
        box.line.width = Pt(2)
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(0.8), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(36)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.7), Inches(box_width - 0.6), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.1), Inches(box_width - 0.6), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        
        col += 1
        if col >= 3:
            col = 0
            row += 1
    
    # Slide 4: Agentic AI Architecture Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Understanding Agentic AI Systems"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Central Agent
    center_x, center_y = 4.25, 3.5
    add_shape_with_text(slide, MSO_SHAPE.HEXAGON, center_x, center_y, 1.5, 1.2, 
                       "AI\nAgent", TITLE_COLOR, font_size=18)
    
    # Surrounding components
    components = [
        (2, 1.5, "Foundation\nModel", BLUE_1),
        (6.5, 1.5, "Tools &\nAPIs", ORANGE),
        (2, 5, "Memory\nSystem", PURPLE),
        (6.5, 5, "Environment\n& Tasks", GREEN),
    ]
    
    for x, y, label, color in components:
        add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.5, 0.8, 
                           label, color, font_size=13)
        # Connect to center
        if x < center_x:
            add_connector(slide, x + 1.5, y + 0.4, center_x, center_y + 0.6, color)
        else:
            add_connector(slide, center_x + 1.5, center_y + 0.6, x, y + 0.4, color)
    
    # Capabilities on the right
    cap_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(1.5), Inches(3))
    tf = cap_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Capabilities:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    for cap in ["🎯 Reason", "📋 Plan", "🔧 Act", "🔄 Learn", "🤝 Collaborate"]:
        p = tf.add_paragraph()
        p.text = cap
        p.font.size = Pt(12)
        p.space_after = Pt(4)
    
    # Slide 5: Multi-Agent Systems
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Agent Systems & Orchestration"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Orchestrator at top
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.75, 1.2, 2.5, 0.7, 
                       "Agent Orchestrator", TITLE_COLOR, font_size=16)
    
    # Specialized agents below
    agents = [
        (1, 2.5, "Research\nAgent", BLUE_1),
        (3, 2.5, "Code\nAgent", GREEN),
        (5, 2.5, "Analysis\nAgent", PURPLE),
        (7, 2.5, "QA\nAgent", ORANGE),
    ]
    
    for x, y, label, color in agents:
        add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.5, 0.8, 
                           label, color, font_size=13)
        # Connect to orchestrator
        add_connector(slide, 5, 1.9, x + 0.75, y, color)
    
    # Key points
    points_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3.2))
    tf = points_box.text_frame
    
    points = [
        ("Orchestration", "Coordinating specialized agents for complex workflows"),
        ("Collaboration", "Swarm, hierarchical, and workflow-based patterns"),
        ("Specialization", "70% of MAS will have narrow, focused roles by 2027"),
        ("Governance", "Permission boundaries and audit logs"),
        ("Human-in-Loop", "Approval checkpoints for critical decisions"),
    ]
    
    for title, desc in points:
        p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        run = p.runs[0]
        run.text = f"• {title}: {desc}"
        run.font.bold = False
        p.space_after = Pt(8)
    
    # Slide 6: AI Evaluation with Three Pillars Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Evaluation Frameworks"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Shift: From 'Can AI do this?' to 'How well, at what cost, and for whom?'"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Three pillars
    pillars = [
        (1.5, "Performance\nMetrics", ["Accuracy/F1", "BLEU/ROUGE", "Relevance"], RGBColor(244, 67, 54)),
        (4.25, "Safety &\nEthics", ["Bias Detection", "Toxicity Check", "Fairness"], GREEN),
        (7, "Production\nQuality", ["Hallucination", "Latency/Cost", "Compliance"], BLUE_1),
    ]
    
    for x, title, metrics, color in pillars:
        # Pillar base
        add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, x, 2, 2, 0.8, 
                           title, color, font_size=15)
        
        # Metrics below
        metric_y = 3.2
        for i, metric in enumerate(metrics):
            metric_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                               Inches(x + 0.2), Inches(metric_y + i * 0.7), 
                                               Inches(1.6), Inches(0.5))
            metric_box.fill.solid()
            metric_box.fill.fore_color.rgb = color
            metric_box.fill.transparency = 0.5
            metric_box.line.color.rgb = color
            
            tf = metric_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Foundation
    found_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                       Inches(1), Inches(5.5), Inches(8), Inches(0.6))
    found_box.fill.solid()
    found_box.fill.fore_color.rgb = TITLE_COLOR
    found_box.line.color.rgb = TITLE_COLOR
    
    tf = found_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Continuous Monitoring & Traceability"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Slide 7: Evaluation Platforms (keeping table from original)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Leading Evaluation Platforms (2026)"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Add table
    rows, cols = 8, 3
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(4)
    table.columns[2].width = Inches(3)
    
    # Header row
    headers = ["Platform", "Strengths", "Best For"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    data = [
        ["DeepEval", "Developer-focused, RAG metrics", "Development & Testing"],
        ["Galileo AI", "Hallucination detection", "Production GenAI"],
        ["Arize", "ML observability, drift", "Enterprise Monitoring"],
        ["Patronus AI", "Rubric-based scoring", "Structured Evaluation"],
        ["MLflow", "Experiment tracking", "Custom Workflows"],
        ["RAGAS", "RAG-specific evaluation", "RAG Systems"],
        ["Braintrust", "Dev workflow integration", "End-to-End Platform"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(13)
    
    # Slide 8: Strands SDK with Architecture Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Strands SDK Framework"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.3))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Open-source, code-first Python framework by AWS"
    p.font.size = Pt(14)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # Architecture diagram
    # Application
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1, 1.8, 2, 0.7, 
                       "Your Application", RGBColor(255, 153, 0), font_size=14)
    
    # Arrow
    add_connector(slide, 3, 2.15, 3.8, 2.15)
    
    # Strands Agent (center)
    add_shape_with_text(slide, MSO_SHAPE.HEXAGON, 3.8, 1.8, 2.4, 1, 
                       "Strands Agent", RGBColor(255, 153, 0), font_size=16)
    
    # Components below
    comps = [
        (1.5, 3.2, "Agent Loop", RGBColor(20, 110, 180)),
        (3.5, 3.2, "Tool Registry", RGBColor(35, 47, 62)),
        (5.5, 3.2, "Memory", RGBColor(236, 114, 17)),
    ]
    
    for x, y, label, color in comps:
        add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.5, 0.6, 
                           label, color, font_size=12)
        add_connector(slide, 5, 2.8, x + 0.75, y)
    
    # Model providers
    providers = [
        (1, 4.3, "AWS\nBedrock", RGBColor(255, 153, 0)),
        (2.8, 4.3, "OpenAI", RGBColor(16, 163, 127)),
        (4.6, 4.3, "Anthropic", RGBColor(217, 119, 87)),
    ]
    
    for x, y, label, color in providers:
        add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.4, 0.6, 
                           label, color, font_size=11)
        add_connector(slide, 2.25, 3.8, x + 0.7, y)
    
    # Features on the right
    feat_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(2.2), Inches(3.5))
    tf = feat_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Features"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    for feat in ["🎯 Simplicity", "🔄 Model Agnostic", "🏗️ Production-Ready", "🤖 Multi-Agent", "⚡ AWS Native"]:
        p = tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(12)
        p.space_after = Pt(6)
    
    # Use cases at bottom
    uc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(1.8))
    tf = uc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Use Cases: "
    p.font.size = Pt(14)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Autonomous incident resolution (SRE)  • Multi-step workflow automation  • Research and analysis agents"
    p.font.size = Pt(12)
    
    # Continue with remaining slides...
    # Slide 9: AWS AgentCore
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AWS Bedrock AgentCore Platform"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Platform diagram
    # Core platform
    platform_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(2), Inches(1.5), Inches(6), Inches(3))
    platform_box.fill.solid()
    platform_box.fill.fore_color.rgb = RGBColor(255, 153, 0)
    platform_box.fill.transparency = 0.1
    platform_box.line.color.rgb = RGBColor(255, 153, 0)
    platform_box.line.width = Pt(3)
    
    # Title inside
    plat_title = slide.shapes.add_textbox(Inches(2.5), Inches(1.7), Inches(5), Inches(0.4))
    tf = plat_title.text_frame
    p = tf.paragraphs[0]
    p.text = "AWS Bedrock AgentCore Platform"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 153, 0)
    p.alignment = PP_ALIGN.CENTER
    
    # Services inside
    services = [
        (2.5, 2.3, "Runtime", RGBColor(20, 110, 180)),
        (4.5, 2.3, "Memory", RGBColor(236, 114, 17)),
        (6.5, 2.3, "Gateway", RGBColor(82, 127, 255)),
        (2.5, 3.2, "Identity", RGBColor(76, 175, 80)),
        (4.5, 3.2, "Observability", RGBColor(156, 39, 176)),
        (6.5, 3.2, "Policy", RGBColor(244, 67, 54)),
    ]
    
    for x, y, label, color in services:
        add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.5, 0.5, 
                           label, color, font_size=11)
    
    # Frameworks above
    fw_label = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(1.2), Inches(0.3))
    tf = fw_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Frameworks:"
    p.font.size = Pt(11)
    p.font.bold = True
    
    frameworks = ["Strands", "LangGraph", "CrewAI", "LlamaIndex"]
    for i, fw in enumerate(frameworks):
        fw_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                        Inches(0.5 + i * 1.3), Inches(1.25), Inches(1.1), Inches(0.4))
        fw_box.fill.solid()
        fw_box.fill.fore_color.rgb = RGBColor(35, 47, 62)
        fw_box.line.color.rgb = RGBColor(35, 47, 62)
        
        tf = fw_box.text_frame
        p = tf.paragraphs[0]
        p.text = fw
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Models below
    model_label = slide.shapes.add_textbox(Inches(8.5), Inches(2), Inches(1), Inches(0.3))
    tf = model_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Models:"
    p.font.size = Pt(11)
    p.font.bold = True
    
    models = ["Claude", "Nova", "Llama", "Mistral"]
    for i, model in enumerate(models):
        model_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                           Inches(8.5), Inches(2.35 + i * 0.5), Inches(1), Inches(0.35))
        model_box.fill.solid()
        model_box.fill.fore_color.rgb = RGBColor(82, 127, 255)
        model_box.line.color.rgb = RGBColor(82, 127, 255)
        
        tf = model_box.text_frame
        p = tf.paragraphs[0]
        p.text = model
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Key capabilities at bottom
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    tf = cap_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Capabilities"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    caps = [
        "• Runtime: Serverless, scalable execution with session isolation",
        "• Memory: Short-term and long-term context management",
        "• Observability: Production monitoring and quality tracking"
    ]
    
    for cap in caps:
        p = tf.add_paragraph()
        p.text = cap
        p.font.size = Pt(12)
        p.space_after = Pt(6)
    
    # Slide 10: Framework Comparison (table from original)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Agent Framework Comparison (2026)"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Add comparison table
    rows, cols = 5, 5
    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(9.4)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2)
    table.columns[4].width = Inches(1.4)
    
    # Header row
    headers = ["Framework", "Philosophy", "Architecture", "Best For", "Learning Curve"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    data = [
        ["LangChain/\nLangGraph", "Modular, flexible", "Graph-based state machines", "Complex enterprise RAG", "Moderate-High"],
        ["CrewAI", "Role-based teams", "Hierarchical/sequential", "Content creation, SOPs", "Easy"],
        ["AutoGen", "Multi-agent conversation", "Dialogue orchestration", "Code gen, research", "Moderate"],
        ["Strands SDK", "Code-first simplicity", "Lightweight agent loop", "AWS-native production", "Easy-Moderate"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Slide 11: LLM Models Comparison (table from original)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Latest LLM Models Comparison (2026)"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Add model comparison table
    rows, cols = 7, 5
    left = Inches(0.2)
    top = Inches(1.2)
    width = Inches(9.6)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    headers = ["Vendor", "Model", "Strengths", "Context Window", "Key Features"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    data = [
        ["OpenAI", "GPT-5.2", "Reasoning, coding", "Extended", "Super-assistant"],
        ["Anthropic", "Claude Opus 4.6", "Agentic coding, safety", "1M tokens", "Computer use"],
        ["Google", "Gemini 3 Pro", "Multimodal, integration", "2M tokens", "Deep Research"],
        ["AWS", "Amazon Nova", "Enterprise, cost-effective", "Variable", "Bedrock-native"],
        ["Meta", "Llama 3+", "Open-source", "Variable", "On-prem deployment"],
        ["Microsoft", "Copilot (GPT-4)", "Productivity", "Extended", "M365 integration"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Slide 12: Cloud Services with Logos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cloud AI Services Overview"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Three cloud providers
    providers_data = [
        (0.5, "Amazon Web Services", RGBColor(255, 153, 0), 
         ["Bedrock - Multi-model platform", "SageMaker - Custom ML", "Amazon Q - AI assistant", "AgentCore - Agentic platform"]),
        (3.5, "Microsoft Azure", RGBColor(0, 120, 212), 
         ["Azure OpenAI - Enterprise GPT-4", "Azure ML - Custom models", "Copilot - M365 integration", "Semantic Kernel - Framework"]),
        (6.5, "Google Cloud", RGBColor(66, 133, 244), 
         ["Vertex AI - ML platform", "Gemini API - Multimodal", "Duet AI - Workspace", "Agent Builder - Workflows"]),
    ]
    
    for x, name, color, services in providers_data:
        # Provider box
        prov_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(x), Inches(1.3), Inches(2.8), Inches(5.2))
        prov_box.fill.solid()
        prov_box.fill.fore_color.rgb = color
        prov_box.fill.transparency = 0.1
        prov_box.line.color.rgb = color
        prov_box.line.width = Pt(2)
        
        # Provider name
        name_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.5), Inches(2.4), Inches(0.5))
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # Services
        serv_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.3), Inches(2.4), Inches(3.8))
        tf = serv_box.text_frame
        tf.word_wrap = True
        
        for service in services:
            p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
            p.text = f"• {service}"
            p.font.size = Pt(11)
            p.space_after = Pt(8)
    
    # Slide 13: Use Cases with Icons
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Real-World Enterprise Use Cases"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Use case boxes
    use_cases = [
        ("💬", "Customer Service", ["Intelligent ticket routing", "24/7 automated support"], BLUE_1),
        ("🔧", "DevOps & SRE", ["Incident auto-remediation", "Log analysis & RCA"], GREEN),
        ("📊", "Data & Analytics", ["Pipeline monitoring", "NL-to-SQL generation"], PURPLE),
        ("💻", "Software Dev", ["Code generation & review", "Bug detection & fixing"], ORANGE),
        ("💰", "Finance", ["Compliance checking", "Fraud detection"], RGBColor(244, 67, 54)),
    ]
    
    col = 0
    row = 0
    start_x = 0.5
    start_y = 1.3
    box_width = 2.8
    box_height = 1.6
    gap_x = 0.4
    gap_y = 0.3
    
    for icon, title, details, color in use_cases:
        x = start_x + col * (box_width + gap_x)
        y = start_y + row * (box_height + gap_y)
        
        # Box
        uc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                        Inches(x), Inches(y), Inches(box_width), Inches(box_height))
        uc_box.fill.solid()
        uc_box.fill.fore_color.rgb = color
        uc_box.fill.transparency = 0.2
        uc_box.line.color.rgb = color
        uc_box.line.width = Pt(2)
        
        # Icon and title
        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(box_width - 0.4), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon} {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Details
        det_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(box_width - 0.4), Inches(0.8))
        tf = det_box.text_frame
        tf.word_wrap = True
        
        for detail in details:
            p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
            p.text = f"• {detail}"
            p.font.size = Pt(11)
            p.space_after = Pt(4)
        
        col += 1
        if col >= 3:
            col = 0
            row += 1
    
    # Success metrics at bottom
    metrics_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                         Inches(0.5), Inches(5.5), Inches(9), Inches(1.3))
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = TITLE_COLOR
    metrics_box.fill.transparency = 0.1
    metrics_box.line.color.rgb = TITLE_COLOR
    
    met_text = slide.shapes.add_textbox(Inches(1), Inches(5.7), Inches(8), Inches(0.9))
    tf = met_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Success Metrics: 📊 70% ↓ resolution time  💰 40% cost savings  ⚡ 90% automation rate  😊 35% ↑ satisfaction"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 14: Best Practices
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Best Practices & Recommendations"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Best practices as numbered steps
    practices = [
        ("1", "Start with Clear Scope", "Define tasks, success criteria, oversight", BLUE_1),
        ("2", "Choose Right Framework", "Match framework to use case requirements", GREEN),
        ("3", "Robust Evaluation", "Testing, monitoring, safety checks", PURPLE),
        ("4", "Governance & Security", "Permissions, audit logs, approvals", ORANGE),
        ("5", "Iterate on Metrics", "Track success, analyze failures, optimize", RGBColor(244, 67, 54)),
    ]
    
    for i, (num, title, desc, color) in enumerate(practices):
        y = 1.3 + i * 1.1
        
        # Number circle
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.5), Inches(0.5))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = color
        num_circle.line.color.rgb = color
        
        tf = num_circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Title and description
        text_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(7.8), Inches(0.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
    
    # Slide 15: Summary & Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Summary & Next Steps"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Key Takeaways
    take_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.5), Inches(3))
    tf = take_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Takeaways"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    takeaways = [
        "✅ Agents: experimental → enterprise-critical",
        "✅ Multi-agent orchestration is the future",
        "✅ Evaluation & observability essential",
        "✅ Choose frameworks by use case",
        "✅ Security & governance from day one"
    ]
    
    for ta in takeaways:
        p = tf.add_paragraph()
        p.text = ta
        p.font.size = Pt(13)
        p.space_after = Pt(8)
    
    # Next Steps
    next_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.2), Inches(3))
    tf = next_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Recommended Next Steps"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    steps = [
        "1. Pilot Project - Well-defined use case",
        "2. Framework Selection - Match to needs",
        "3. Evaluation Setup - Monitor day one",
        "4. Team Training - Upskill team",
        "5. Production Deployment - Iterate"
    ]
    
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(13)
        p.space_after = Pt(8)
    
    # Resources at bottom
    res_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
    tf = res_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Resources & Documentation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    resources = [
        "Strands SDK: strandsagents.com",
        "AWS AgentCore: docs.aws.amazon.com/bedrock",
        "LangChain: langchain.com  •  CrewAI: crewai.com"
    ]
    
    for res in resources:
        p = tf.add_paragraph()
        p.text = res
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)
    
    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('e:/Antigravity/sreagent/AI_GenAI_Agents_Presentation.pptx')
    print("✅ Presentation with diagrams created successfully!")
    print("📊 15 slides with visual diagrams for each topic")
    return prs

if __name__ == "__main__":
    create_presentation()
